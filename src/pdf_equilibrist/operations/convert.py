"""
operations/convert.py — Conversion de et vers PDF
==================================================
Ce module regroupe toutes les conversions de format impliquant des PDF :
PDF → image, PDF → Word/Excel/PowerPoint, image → PDF, Office → PDF.

Architecture de détection des moteurs
--------------------------------------
La conversion Office → PDF nécessite un moteur externe installé sur le poste.
Deux moteurs sont supportés, par ordre de priorité :

1. **Microsoft Office** (via ``docx2pdf``) :
   Détecté par la présence de clés dans le registre Windows
   ``HKLM\\SOFTWARE\\Microsoft\\Windows\\CurrentVersion\\App Paths\\WINWORD.EXE``.
   Utilise l'automation COM de Word/Excel/PowerPoint — conversion haute fidélité.

2. **LibreOffice** (via ``subprocess soffice --headless``) :
   Détecté par présence de l'exe dans les chemins standards ou dans le PATH.
   Fallback multiplateforme (Windows, Linux, macOS).

Dépendances externes
---------------------
- ``pdf2docx``   : conversion PDF → Word (analyse la structure de mise en page)
- ``pdfplumber`` : extraction de tableaux pour PDF → Excel
- ``openpyxl``   : écriture du fichier .xlsx
- ``python-pptx``: création des slides pour PDF → PowerPoint
- ``docx2pdf``   : automation COM pour Office → PDF (MS Office requis)
- ``Pillow``     : manipulation d'images
- ``PyMuPDF``    : rastérisation des pages PDF
"""
from pathlib import Path
import fitz
import io
import subprocess
import shutil


# ── Détection moteur Office → PDF ────────────────────────────────────────────

def detect_office_engine() -> str | None:
    """
    Détecte quel moteur de conversion Office → PDF est disponible sur le poste.

    Ordre de priorité : MS Office d'abord (meilleure fidélité), LibreOffice ensuite.

    Détection MS Office
    ~~~~~~~~~~~~~~~~~~~
    Recherche les clés de registre Windows ``App Paths`` pour les exécutables
    Word, Excel et PowerPoint. Une seule correspondance suffit pour valider
    la présence de la suite Office.

    Détection LibreOffice
    ~~~~~~~~~~~~~~~~~~~~~
    Vérifie les chemins d'installation standards sur Windows, Linux et macOS,
    puis tente ``shutil.which("soffice")`` pour les installations non standard.

    Returns
    -------
    str | None
        ``'msoffice'``    — Microsoft Office détecté via registre Windows.
        ``'libreoffice'`` — LibreOffice détecté dans le système de fichiers.
        ``None``          — Aucun moteur trouvé.
    """
    # 1. MS Office : vérification via le registre Windows (HKLM)
    try:
        import winreg
        for app in ["WINWORD.EXE", "EXCEL.EXE", "POWERPNT.EXE"]:
            try:
                winreg.OpenKey(
                    winreg.HKEY_LOCAL_MACHINE,
                    rf"SOFTWARE\Microsoft\Windows\CurrentVersion\App Paths\{app}"
                )
                return "msoffice"   # au moins un exe Office trouvé → suffisant
            except OSError:
                continue            # cet exe absent, essayer le suivant
    except ImportError:
        pass   # winreg non disponible (Linux/macOS) → passer à LibreOffice

    # 2. LibreOffice : chemins d'installation courants
    lo_candidates = [
        r"C:\Program Files\LibreOffice\program\soffice.exe",
        r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
        "/usr/bin/soffice",
        "/usr/local/bin/soffice",
        "/Applications/LibreOffice.app/Contents/MacOS/soffice",
    ]
    for p in lo_candidates:
        if Path(p).exists():
            return "libreoffice"

    # Fallback : chercher soffice dans le PATH système
    if shutil.which("soffice"):
        return "libreoffice"

    return None   # aucun moteur disponible


def office_to_pdf(input_path: Path, output_path: Path) -> Path:
    """
    Convertit un fichier Office (.docx, .xlsx, .pptx, .doc, .odt…) en PDF.

    Sélectionne automatiquement le moteur disponible via ``detect_office_engine()``.

    MS Office (via docx2pdf)
    ~~~~~~~~~~~~~~~~~~~~~~~~~
    ``docx2pdf`` utilise l'automation COM Windows pour piloter Word/Excel/PowerPoint
    silencieusement. La conversion est haute fidélité (rendu identique à l'impression).
    Nécessite que l'application Office correspondante soit installée et activée.

    LibreOffice (via subprocess)
    ~~~~~~~~~~~~~~~~~~~~~~~~~~~~
    Lance ``soffice --headless --convert-to pdf`` en sous-processus.
    LibreOffice génère le PDF dans le dossier de sortie avec le nom du fichier source.
    Un renommage est effectué si le nom généré diffère de ``output_path``.

    Parameters
    ----------
    input_path : Path
        Fichier Office source à convertir.
    output_path : Path
        Chemin du PDF de sortie désiré.

    Returns
    -------
    Path
        Chemin du PDF généré (= ``output_path`` si succès).

    Raises
    ------
    RuntimeError
        Si la conversion échoue ou si aucun moteur n'est disponible.
    """
    engine = detect_office_engine()

    if engine == "msoffice":
        try:
            import docx2pdf
            docx2pdf.convert(str(input_path), str(output_path))
            return output_path
        except Exception as e:
            raise RuntimeError(f"Conversion MS Office échouée : {e}") from e

    if engine == "libreoffice":
        # Trouver le chemin exact de soffice
        lo_exe = (
            r"C:\Program Files\LibreOffice\program\soffice.exe"
            if Path(r"C:\Program Files\LibreOffice\program\soffice.exe").exists()
            else shutil.which("soffice")
        )
        result = subprocess.run(
            [
                lo_exe,
                "--headless",               # pas d'interface graphique
                "--convert-to", "pdf",      # format de sortie
                "--outdir", str(output_path.parent),  # dossier de destination
                str(input_path),
            ],
            capture_output=True,
            text=True,
            timeout=120,   # 2 minutes max pour les gros fichiers
        )
        if result.returncode != 0:
            raise RuntimeError(f"LibreOffice a échoué :\n{result.stderr}")

        # LibreOffice nomme le PDF d'après le fichier source : renommer si besoin
        generated = output_path.parent / (input_path.stem + ".pdf")
        if generated != output_path and generated.exists():
            generated.rename(output_path)
        return output_path

    raise RuntimeError(
        "Aucun moteur de conversion trouvé.\n"
        "Installez Microsoft Office ou LibreOffice."
    )


# ── PDF → Formats externes ────────────────────────────────────────────────────

def to_images(
    doc: fitz.Document,
    output_dir: Path,
    zoom: float = 2.0,
    fmt: str = "png",
) -> list[Path]:
    """
    Rastérise toutes les pages du PDF en fichiers image.

    Chaque page produit un fichier ``page_1.png``, ``page_2.png``, etc.
    dans le dossier de sortie.

    Parameters
    ----------
    doc : fitz.Document
        Document PDF source.
    output_dir : Path
        Dossier de destination (créé si absent).
    zoom : float
        Facteur de résolution. 1.0 = 72 dpi (natif PDF), 2.0 = 144 dpi (défaut),
        4.0 = 288 dpi (impression haute qualité). Un zoom élevé produit des
        fichiers plus lourds mais des images plus nettes.
    fmt : str
        Format image de sortie : ``"png"``, ``"jpg"``, ``"bmp"``, ``"tiff"``.

    Returns
    -------
    list[Path]
        Chemins des fichiers image créés, dans l'ordre des pages.
    """
    output_dir.mkdir(parents=True, exist_ok=True)
    mat = fitz.Matrix(zoom, zoom)
    paths = []
    for i, page in enumerate(doc):
        pix = page.get_pixmap(matrix=mat, alpha=False)
        p = output_dir / f"page_{i + 1}.{fmt}"
        pix.save(str(p))
        paths.append(p)
    return paths


def to_word(doc: fitz.Document, src_path: Path, output_path: Path) -> Path:
    """
    Convertit un PDF en document Word (.docx) via pdf2docx.

    ``pdf2docx`` analyse la mise en page du PDF (blocs texte, images, tableaux)
    et tente de reproduire la structure dans un document Word éditable.
    La fidélité dépend de la complexité du PDF original.

    Note : ``pdf2docx`` a besoin du **chemin fichier** (pas du ``fitz.Document``
    en mémoire) car il ouvre le PDF lui-même en interne.

    Parameters
    ----------
    doc : fitz.Document
        Document PyMuPDF (non utilisé directement, présent pour cohérence API).
    src_path : Path
        Chemin du fichier PDF source sur le disque.
    output_path : Path
        Chemin du fichier .docx de sortie.

    Returns
    -------
    Path
        Chemin du fichier .docx généré.
    """
    from pdf2docx import Converter
    cv = Converter(str(src_path))
    cv.convert(str(output_path))
    cv.close()
    return output_path


def to_excel(doc: fitz.Document, output_path: Path) -> Path:
    """
    Extrait les tableaux du PDF et les écrit dans un fichier Excel (.xlsx).

    Utilise ``pdfplumber`` pour détecter et extraire les tableaux page par page.
    Chaque page produit une feuille Excel distincte (``Page 1``, ``Page 2``, …).
    Si une page ne contient pas de tableau, sa feuille sera vide.

    Limitation : la détection de tableaux fonctionne sur les PDF avec des lignes
    de bordure visibles. Les tableaux sans bordure (espacement seul) peuvent
    ne pas être détectés correctement.

    Parameters
    ----------
    doc : fitz.Document
        Document PyMuPDF (converti en bytes pour pdfplumber via ``io.BytesIO``).
    output_path : Path
        Chemin du fichier .xlsx de sortie.

    Returns
    -------
    Path
        Chemin du fichier Excel généré.
    """
    import pdfplumber
    import openpyxl

    # pdfplumber ne peut pas lire un fitz.Document directement →
    # on sérialise le doc en bytes en mémoire pour l'ouvrir via BytesIO
    buf = io.BytesIO(doc.tobytes())
    wb = openpyxl.Workbook()
    first = True

    with pdfplumber.open(buf) as pdf:
        for i, page in enumerate(pdf.pages):
            # Créer ou réutiliser la feuille active (openpyxl crée "Sheet" par défaut)
            ws = wb.active if first else wb.create_sheet(f"Page {i + 1}")
            first = False
            ws.title = f"Page {i + 1}"

            # extract_tables() retourne une liste de tableaux,
            # chaque tableau est une liste de lignes (list of list of str)
            tables = page.extract_tables()
            for table in tables:
                for row in table:
                    # Remplacer None (cellule vide détectée) par chaîne vide
                    ws.append([cell or "" for cell in row])
                ws.append([])   # ligne vide entre deux tableaux

    wb.save(str(output_path))
    return output_path


def to_powerpoint(
    doc: fitz.Document,
    output_path: Path,
    zoom: float = 1.5,
) -> Path:
    """
    Convertit chaque page du PDF en slide PowerPoint (.pptx).

    Chaque page est rastérisée en image PNG puis insérée comme image plein-slide
    dans une présentation PowerPoint. Le format de la slide est adapté aux
    dimensions exactes de chaque page PDF (portrait, paysage, formats non standard).

    Note : la conversion est « image de page », pas une conversion structurelle.
    Le texte des slides ne sera pas éditable dans PowerPoint.

    Parameters
    ----------
    doc : fitz.Document
        Document PDF source.
    output_path : Path
        Chemin du fichier .pptx de sortie.
    zoom : float
        Résolution de rastérisation (1.5 = 108 dpi, bon compromis taille/qualité).

    Returns
    -------
    Path
        Chemin du fichier .pptx généré.
    """
    from pptx import Presentation
    from pptx.util import Inches

    mat = fitz.Matrix(zoom, zoom)
    prs = Presentation()
    blank_layout = prs.slide_layouts[6]   # layout "Vierge" sans placeholder

    for page in doc:
        # Rastériser la page en PNG en mémoire
        pix = page.get_pixmap(matrix=mat, alpha=False)
        buf = io.BytesIO(pix.tobytes("png"))

        # Ajouter un slide et ajuster ses dimensions à la page PDF
        # Les dimensions PDF sont en points (1 pt = 1/72 pouce)
        width_in  = page.rect.width  / 72   # largeur en pouces
        height_in = page.rect.height / 72   # hauteur en pouces
        prs.slide_width  = Inches(width_in)
        prs.slide_height = Inches(height_in)

        slide = prs.slides.add_slide(blank_layout)
        # Insérer l'image en plein slide (position 0,0, taille = slide entière)
        slide.shapes.add_picture(buf, 0, 0, Inches(width_in), Inches(height_in))

    prs.save(str(output_path))
    return output_path


# ── Formats externes → PDF ────────────────────────────────────────────────────

def image_to_pdf(image_paths: list[Path], output_path: Path) -> Path:
    """
    Assemble plusieurs images en un seul fichier PDF.

    Chaque image devient une page PDF de format A4 (595×842 points).
    L'image est étirée pour remplir la page entière.

    Formats acceptés : PNG, JPEG, BMP, TIFF, GIF — tout ce que PyMuPDF supporte.

    Parameters
    ----------
    image_paths : list[Path]
        Liste ordonnée des images à inclure. Chaque image → une page PDF.
    output_path : Path
        Chemin du fichier PDF de sortie.

    Returns
    -------
    Path
        Chemin du PDF généré.
    """
    doc = fitz.open()
    for img_path in image_paths:
        # Créer une page A4 et y insérer l'image en plein cadre
        page = doc.new_page(width=595, height=842)   # A4 en points PDF
        page.insert_image(page.rect, filename=str(img_path))
    doc.save(str(output_path))
    doc.close()
    return output_path


# ── PDF → Markdown ────────────────────────────────────────────────────────────

def to_markdown(doc: fitz.Document, output_path: Path) -> Path:
    """
    Convertit un PDF en fichier Markdown (.md).

    Algorithme
    ----------
    1. Collecter toutes les tailles de police du document pour déterminer la
       taille « corps » (mode statistique) et les niveaux de titres.
    2. Pour chaque page, détecter les tableaux via pdfplumber (zones à exclure
       du flux texte normal).
    3. Extraire les blocs texte via PyMuPDF ``get_text("dict")``, en sautant
       les zones couvertes par des tableaux.
    4. Traduire chaque span en Markdown selon la taille de police et les flags
       (gras, italique).
    5. Insérer les tableaux pdfplumber au bon endroit (position Y dans la page).

    Qualité selon le type de PDF
    -----------------------------
    - Document Office exporté en PDF → résultat fidèle
    - Rapport structuré → correct
    - Plan technique / CAD → médiocre (layout complexe, textes épars)
    - PDF scanné (image) → lève ``ValueError``

    Parameters
    ----------
    doc : fitz.Document
        Document PDF source ouvert.
    output_path : Path
        Chemin du fichier .md de sortie.

    Returns
    -------
    Path
        Chemin du fichier .md généré.

    Raises
    ------
    ValueError
        Si le document est un scan (aucune couche texte extractible).
    """
    import statistics
    import pdfplumber

    # ── Passe 1 : collecter toutes les tailles pour déduire le corps ──────────
    all_sizes: list[float] = []
    for page in doc:
        raw = page.get_text("dict", flags=fitz.TEXT_PRESERVE_WHITESPACE)
        for block in raw.get("blocks", []):
            if block.get("type") != 0:
                continue
            for line in block.get("lines", []):
                for span in line.get("spans", []):
                    t = span.get("text", "").strip()
                    if t:
                        all_sizes.append(round(span.get("size", 12.0), 1))

    if not all_sizes:
        raise ValueError(
            "Aucun texte extractible — ce PDF est probablement un scan.\n"
            "La conversion PDF → Markdown nécessite une couche texte."
        )

    body_size = statistics.mode(all_sizes)   # taille la plus fréquente = corps

    def _heading_level(size: float) -> int:
        """Retourne 0 (corps) ou 1/2/3 (titre H1/H2/H3) selon la taille."""
        ratio = size / body_size
        if ratio >= 1.8:
            return 1
        if ratio >= 1.4:
            return 2
        if ratio >= 1.15:
            return 3
        return 0

    def _inline(text: str, flags: int) -> str:
        """Applique le formatage inline gras/italique."""
        bold   = bool(flags & 0b10000)   # bit 4 = bold dans PyMuPDF
        italic = bool(flags & 0b01)      # bit 0 = italic
        text = text.replace("\\", "\\\\")
        if bold and italic:
            return f"***{text}***"
        if bold:
            return f"**{text}**"
        if italic:
            return f"*{text}*"
        return text

    def _table_to_md(table: list[list]) -> str:
        """Convertit une table pdfplumber (liste de listes) en Markdown."""
        if not table:
            return ""
        rows = []
        for row in table:
            cells = [str(c).strip().replace("\n", " ") if c else "" for c in row]
            rows.append("| " + " | ".join(cells) + " |")
        if len(rows) >= 1:
            # Ligne de séparation après le header
            sep = "| " + " | ".join(["---"] * len(table[0])) + " |"
            rows.insert(1, sep)
        return "\n".join(rows)

    # ── Passe 2 : extraction page par page ────────────────────────────────────
    md_lines: list[str] = []
    pdf_path = doc.name   # chemin du fichier source pour pdfplumber

    with pdfplumber.open(pdf_path) as plumber_doc:
        for page_idx, page in enumerate(doc):
            plumber_page = plumber_doc.pages[page_idx]

            # Détecter les tableaux et leur position Y (coord pdfplumber = top-based)
            # pdfplumber utilise l'origine en haut-gauche, PyMuPDF aussi.
            pl_tables   = plumber_page.find_tables()
            table_bboxes: list[tuple] = []   # (x0, y0, x1, y1) en points PDF
            table_md:    dict[float, str] = {}  # y0 → markdown du tableau

            for tbl in pl_tables:
                data = tbl.extract()
                if not data:
                    continue
                bbox = tbl.bbox   # (x0, top, x1, bottom) dans le repère pdfplumber
                table_bboxes.append(bbox)
                table_md[bbox[1]] = _table_to_md(data)

            def _in_table(bbox_span: tuple) -> bool:
                """Retourne True si ce span est dans une zone tableau."""
                sx0, sy0, sx1, sy1 = bbox_span
                for tx0, ty0, tx1, ty1 in table_bboxes:
                    overlap_x = sx0 < tx1 and sx1 > tx0
                    overlap_y = sy0 < ty1 and sy1 > ty0
                    if overlap_x and overlap_y:
                        return True
                return False

            # Construire la liste des éléments de la page triés par Y
            # Chaque élément : (y_position, type, content)
            # type = "text" | "table"
            elements: list[tuple[float, str, str]] = []

            # Ajouter les tableaux
            for ty0, md_table in table_md.items():
                elements.append((ty0, "table", md_table))

            # Extraire les blocs texte hors tableaux
            raw = page.get_text("dict", flags=fitz.TEXT_PRESERVE_WHITESPACE)
            for block in raw.get("blocks", []):
                if block.get("type") != 0:
                    continue

                for line in block.get("lines", []):
                    line_parts: list[str] = []
                    line_y     = line["bbox"][1]
                    line_size  = 0.0

                    for span in line.get("spans", []):
                        text = span.get("text", "")
                        if not text.strip():
                            continue
                        if _in_table(span["bbox"]):
                            continue

                        s    = span.get("size", body_size)
                        flags = span.get("flags", 0)
                        line_size = max(line_size, s)
                        line_parts.append(_inline(text, flags))

                    if not line_parts:
                        continue

                    line_text = "".join(line_parts).strip()
                    if not line_text:
                        continue

                    level = _heading_level(line_size)
                    if level:
                        formatted = "#" * level + " " + line_text
                    else:
                        formatted = line_text

                    elements.append((line_y, "text", formatted))

            # Trier par position Y croissante (haut → bas)
            elements.sort(key=lambda e: e[0])

            # Assembler les lignes de la page
            prev_y    = -1.0
            prev_type = ""
            for y, etype, content in elements:
                # Insérer une ligne vide entre éléments espacés (>1 interligne)
                gap = y - prev_y
                if prev_y >= 0 and gap > body_size * 1.8 and prev_type == "text":
                    md_lines.append("")

                if etype == "table":
                    if md_lines and md_lines[-1] != "":
                        md_lines.append("")
                    md_lines.append(content)
                    md_lines.append("")
                else:
                    md_lines.append(content)

                prev_y    = y
                prev_type = etype

            # Séparateur de page (sauf après la dernière)
            if page_idx < len(doc) - 1:
                md_lines.append("")
                md_lines.append("---")
                md_lines.append("")

    # ── Nettoyage final ───────────────────────────────────────────────────────
    # Supprimer les lignes vides consécutives (max 1 ligne vide d'affilée)
    cleaned: list[str] = []
    prev_blank = False
    for line in md_lines:
        is_blank = not line.strip()
        if is_blank and prev_blank:
            continue
        cleaned.append(line)
        prev_blank = is_blank

    output_path.write_text("\n".join(cleaned), encoding="utf-8")
    return output_path
