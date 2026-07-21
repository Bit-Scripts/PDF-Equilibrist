# -*- coding: utf-8 -*-
"""Génère la V2 de la présentation PDF Equilibrist."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# ── Palette ──────────────────────────────────────────────────────────────────
DARK      = RGBColor(0x1E, 0x1E, 0x1E)
DARK2     = RGBColor(0x2D, 0x2D, 0x2D)
GREEN     = RGBColor(0x6B, 0xBF, 0x4E)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
OFFWHITE  = RGBColor(0xF5, 0xF5, 0xF5)
MIDGRAY   = RGBColor(0x8A, 0x8A, 0x8A)
DARKTEXT  = RGBColor(0x22, 0x22, 0x22)
LIGHTBG   = RGBColor(0xF8, 0xF8, 0xF8)
GREENBG   = RGBColor(0xEB, 0xF7, 0xE5)
CARDBG    = RGBColor(0xF0, 0xF0, 0xF0)
GREENBORDER = RGBColor(0x55, 0xA0, 0x3A)

W = Inches(13.33)   # 960 pt → 16:9
H = Inches(7.5)     # 540 pt

# ── Helpers ───────────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def blank_slide(prs):
    blank_layout = prs.slide_layouts[6]   # Vide
    return prs.slides.add_slide(blank_layout)


def fill_solid(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def no_line(shape):
    shape.line.fill.background()


def add_rect(slide, x, y, w, h, fill=None, line=None, line_w=Pt(1)):
    from pptx.util import Pt
    shape = slide.shapes.add_shape(1, x, y, w, h)   # MSO_SHAPE_TYPE.RECTANGLE = 1
    if fill:
        fill_solid(shape, fill)
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = line_w
    else:
        no_line(shape)
    shape.shadow.inherit = False
    return shape


def add_text(slide, text, x, y, w, h,
             size=Pt(14), bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = size
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txb


def add_para(tf, text, size=Pt(13), bold=False, color=DARKTEXT,
             align=PP_ALIGN.LEFT, italic=False, level=0, space_before=Pt(0)):
    from pptx.util import Pt
    p = tf.add_paragraph()
    p.level = level
    p.alignment = align
    p.space_before = space_before
    run = p.add_run()
    run.text = text
    run.font.size   = size
    run.font.bold   = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return p


def header_bar(slide, title, subtitle=None):
    """Barre supérieure foncée avec titre vert."""
    bar_h = Inches(1.15)
    add_rect(slide, 0, 0, W, bar_h, fill=DARK)
    add_rect(slide, 0, bar_h, W, Inches(0.045), fill=GREEN)  # ligne accent

    add_text(slide, title,
             Inches(0.45), Inches(0.15), Inches(10), Inches(0.65),
             size=Pt(28), bold=True, color=WHITE)
    if subtitle:
        add_text(slide, subtitle,
                 Inches(0.45), Inches(0.72), Inches(10), Inches(0.38),
                 size=Pt(14), color=GREEN, bold=False)


def bullet_box(slide, items, x, y, w, h, title=None,
               bg=LIGHTBG, title_color=GREEN, text_color=DARKTEXT,
               text_size=Pt(13), bullet="▸ ", title_size=Pt(15)):
    """Boîte avec titre optionnel + liste à puces."""
    box = add_rect(slide, x, y, w, h, fill=bg, line=CARDBG)
    if title:
        title_h = Inches(0.38)
        title_bg = add_rect(slide, x, y, w, title_h, fill=DARK2)
        add_text(slide, title, x + Inches(0.15), y + Inches(0.05),
                 w - Inches(0.2), title_h - Inches(0.05),
                 size=title_size, bold=True, color=GREEN)
        y += title_h

    txb = slide.shapes.add_textbox(x + Inches(0.18), y + Inches(0.12),
                                   w - Inches(0.25), h - Inches(0.2) - (Inches(0.38) if title else 0))
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = bullet + item
        run.font.size = text_size
        run.font.color.rgb = text_color


# ── Slide 1 – Titre ──────────────────────────────────────────────────────────

def slide_titre(prs):
    s = blank_slide(prs)

    # Fond foncé total
    add_rect(s, 0, 0, W, H, fill=DARK)

    # Bande verte verticale gauche
    add_rect(s, 0, 0, Inches(0.22), H, fill=GREEN)

    # Ligne verte décorative horizontale
    add_rect(s, Inches(0.22), Inches(3.7), W - Inches(0.22), Inches(0.05), fill=GREEN)

    # Titre principal
    add_text(s, "PDF Equilibrist",
             Inches(0.55), Inches(1.4), Inches(11), Inches(1.1),
             size=Pt(54), bold=True, color=WHITE)

    # Sous-titre
    add_text(s, "Éditeur PDF desktop Windows",
             Inches(0.55), Inches(2.6), Inches(9), Inches(0.55),
             size=Pt(26), bold=False, color=GREEN)

    # Descriptif
    add_text(s, "Besoin métier · Fonctionnalités · Architecture · Déploiement",
             Inches(0.55), Inches(3.9), Inches(10), Inches(0.45),
             size=Pt(15), bold=False, color=MIDGRAY)

    # Date / version
    add_text(s, "Juillet 2026  —  v2.0",
             Inches(0.55), Inches(6.8), Inches(5), Inches(0.4),
             size=Pt(12), color=MIDGRAY)


# ── Slide 2 – Aperçu de l'application ───────────────────────────────────────

def slide_apercu(prs, screenshot_path):
    s = blank_slide(prs)
    add_rect(s, 0, 0, W, H, fill=DARK)
    add_rect(s, 0, 0, W, Inches(0.8), fill=DARK2)
    add_rect(s, 0, Inches(0.8), W, Inches(0.04), fill=GREEN)

    add_text(s, "Aperçu de l'application",
             Inches(0.4), Inches(0.1), Inches(10), Inches(0.65),
             size=Pt(22), bold=True, color=WHITE)

    # Screenshot centré
    img_w = Inches(11.8)
    img_h = Inches(6.2)
    img_x = (W - img_w) / 2
    img_y = Inches(0.95)
    if os.path.exists(screenshot_path):
        s.shapes.add_picture(screenshot_path, img_x, img_y, img_w, img_h)


# ── Slide 3 – Le besoin métier ───────────────────────────────────────────────

def slide_besoin(prs):
    s = blank_slide(prs)
    add_rect(s, 0, 0, W, H, fill=OFFWHITE)
    header_bar(s, "1 — Le besoin métier", "Contexte et problématique")

    # 3 blocs côte à côte
    bw = Inches(3.9)
    bh = Inches(4.8)
    by = Inches(1.5)

    # Bloc 1
    bullet_box(s, [
        "Les PDF dominent la diffusion des documents finals",
        "Modification impossible sans logiciel dédié",
        "Outils lourds ou coûteux (Adobe Acrobat, Nitro…)",
        "Dépendance à des solutions cloud non maîtrisées",
    ], Inches(0.3), by, bw, bh, title="Problème actuel",
       bg=WHITE, text_size=Pt(13))

    # Bloc 2
    bullet_box(s, [
        "Corriger ou enrichir un PDF sans re-traitement",
        "Annoter, sécuriser, convertir depuis Windows",
        "Flux de travail local, sans dépendance cloud",
        "Impression fidèle des documents techniques",
    ], Inches(4.72), by, bw, bh, title="Besoin ciblé",
       bg=WHITE, text_size=Pt(13))

    # Bloc 3
    bullet_box(s, [
        "Équipes administratives et de documentation",
        "Services production / coordination de projets",
        "Gestion de documents sensibles ou règlementés",
        "Plans techniques AutoCAD, documents multi-pages",
    ], Inches(9.15), by, bw, bh, title="Qui est concerné",
       bg=WHITE, text_size=Pt(13))


# ── Slide 4 – Fonctionnalités ────────────────────────────────────────────────

def slide_fonctionnalites(prs):
    s = blank_slide(prs)
    add_rect(s, 0, 0, W, H, fill=OFFWHITE)
    header_bar(s, "2 — Fonctionnalités", "6 domaines fonctionnels couverts")

    cols = [
        ("👁  Afficher",    ["Zoom 25–400 %", "Rotation pages", "Miniatures latérales", "Multi-documents (onglets)"]),
        ("✏️  Modifier",    ["Édition de texte inline", "Ajout texte / image", "Filigrane", "Tampon / signature"]),
        ("📋  Annoter",     ["Surlignage", "Barré / Souligné", "Zone de texte libre", "Placement flottant interactif"]),
        ("📄  Pages",       ["Insertion / découpage", "Fusion de PDFs", "Rotation / inversion", "Recadrage pages"]),
        ("🔄  Convertir",   ["PDF → Word / Excel / PPT", "PDF → Images", "Image / Office → PDF", "Traitement par lot"]),
        ("🔒  Protéger",    ["Chiffrement AES-256", "Déchiffrement", "Gestion des permissions", "Impression sécurisée"]),
    ]

    bw = Inches(2.05)
    bh = Inches(4.65)
    by = Inches(1.45)
    gap = Inches(0.06)
    x0 = Inches(0.22)

    for i, (title, items) in enumerate(cols):
        bx = x0 + i * (bw + gap)
        bullet_box(s, items, bx, by, bw, bh, title=title,
                   bg=WHITE, text_size=Pt(11.5))


# ── Slide 5 – Utilisateurs ───────────────────────────────────────────────────

def slide_utilisateurs(prs):
    s = blank_slide(prs)
    add_rect(s, 0, 0, W, H, fill=OFFWHITE)
    header_bar(s, "3 — Utilisateurs et cas d'usage", "Profils prioritaires identifiés")

    rows = [
        ("Équipes administratives",
         "Correction de texte, mise à jour de documents, annotations de relecture, ajout de signatures."),
        ("Services production / documentation",
         "Filigranes, tampons de validation, conversions Office↔PDF, traitement par lot de dossiers."),
        ("Direction / coordination",
         "Chiffrement AES-256, contrôle des permissions, diffusion de documents sensibles ou réglementés."),
        ("Bureaux techniques",
         "Impression de plans AutoCAD (traits fins, hachures), fusion de livrables multi-pages, recadrage."),
    ]

    icons = ["🏢", "📦", "👔", "📐"]
    bh = Inches(1.38)
    bw_label = Inches(3.2)
    bw_text  = Inches(9.45)
    x0 = Inches(0.28)
    y0 = Inches(1.5)
    gap = Inches(0.15)

    for i, (label, text) in enumerate(rows):
        by = y0 + i * (bh + gap)

        # Icône + label (fond vert sombre)
        add_rect(s, x0, by, bw_label, bh, fill=DARK2)
        add_text(s, icons[i] + "  " + label,
                 x0 + Inches(0.15), by + Inches(0.38),
                 bw_label - Inches(0.2), Inches(0.55),
                 size=Pt(14), bold=True, color=GREEN)

        # Texte descriptif
        add_rect(s, x0 + bw_label, by, bw_text, bh, fill=WHITE)
        add_text(s, text,
                 x0 + bw_label + Inches(0.2), by + Inches(0.22),
                 bw_text - Inches(0.3), bh - Inches(0.3),
                 size=Pt(13.5), color=DARKTEXT, wrap=True)


# ── Slide 6 – Architecture et technique ──────────────────────────────────────

def slide_technique(prs):
    s = blank_slide(prs)
    add_rect(s, 0, 0, W, H, fill=OFFWHITE)
    header_bar(s, "4 — Architecture technique", "Stack open-source · Sans droits admin · Windows 10/11")

    # Colonne gauche — Stack
    bullet_box(s, [
        "PyQt6 — interface graphique native Windows",
        "PyMuPDF (fitz) — moteur PDF open-source",
        "PyInstaller — exe autonome (~100 Mo)",
        "Chiffrement AES-256 via PyMuPDF",
        "Impression Win32 GDI — rendu vectoriel haute qualité",
    ], Inches(0.3), Inches(1.5), Inches(5.8), Inches(4.8),
       title="Stack technique (MIT / AGPL-3)", bg=WHITE, text_size=Pt(13))

    # Colonne droite — Déploiement
    bullet_box(s, [
        "Installation NSIS sans droits administrateur (HKCU)",
        "Aucune donnée transmise — 100 % local",
        "Aucune dépendance cloud en fonctionnement",
        "CI/CD : build automatique sur tag Git (GitHub Actions)",
        "Désinstallation propre via Programmes et fonctionnalités",
    ], Inches(6.55), Inches(1.5), Inches(6.5), Inches(2.8),
       title="Déploiement", bg=WHITE, text_size=Pt(13))

    # Dépendances tierces
    bullet_box(s, [
        "MS Office ou LibreOffice — requis pour conversion Office→PDF",
        "Windows 10 / 11 uniquement (Win32 GDI, APIs Qt6)",
    ], Inches(6.55), Inches(4.55), Inches(6.5), Inches(1.75),
       title="Dépendances optionnelles", bg=WHITE, text_size=Pt(13))


# ── Slide 7 – Valeur fonctionnelle ───────────────────────────────────────────

def slide_valeur(prs):
    s = blank_slide(prs)
    add_rect(s, 0, 0, W, H, fill=OFFWHITE)
    header_bar(s, "5 — Valeur fonctionnelle", "Apports concrets pour les équipes")

    cards = [
        ("Productivité",
         "Remplace plusieurs outils disparates par une interface unique, locale et maîtrisée."),
        ("Conformité",
         "Chiffrement AES-256, gestion des permissions et traçabilité des accès aux documents sensibles."),
        ("Qualité d'impression",
         "Rendu fidèle des plans techniques (AutoCAD) : traits fins visibles, hachures filtrées."),
        ("Maintenabilité",
         "Code Python structuré, open-source, versioning Git, CI/CD — évolutions et correctifs facilitées."),
    ]

    bw = Inches(6.2)
    bh = Inches(2.35)
    positions = [
        (Inches(0.25), Inches(1.5)),
        (Inches(6.85), Inches(1.5)),
        (Inches(0.25), Inches(4.05)),
        (Inches(6.85), Inches(4.05)),
    ]
    icons = ["⚡", "🔒", "🖨️", "🔧"]

    for (bx, by), (title, text), icon in zip(positions, cards, icons):
        add_rect(s, bx, by, bw, bh, fill=WHITE, line=CARDBG)
        # Barre latérale verte
        add_rect(s, bx, by, Inches(0.09), bh, fill=GREEN)
        add_text(s, icon + "  " + title,
                 bx + Inches(0.18), by + Inches(0.12),
                 bw - Inches(0.25), Inches(0.5),
                 size=Pt(16), bold=True, color=DARKTEXT)
        add_text(s, text,
                 bx + Inches(0.18), by + Inches(0.65),
                 bw - Inches(0.3), bh - Inches(0.75),
                 size=Pt(13.5), color=MIDGRAY, wrap=True)


# ── Slide 8 – Prochaines étapes ──────────────────────────────────────────────

def slide_etapes(prs):
    s = blank_slide(prs)
    add_rect(s, 0, 0, W, H, fill=DARK)
    add_rect(s, 0, 0, Inches(0.22), H, fill=GREEN)
    add_rect(s, Inches(0.22), Inches(1.4), W - Inches(0.22), Inches(0.04), fill=GREEN)

    add_text(s, "Prochaines étapes",
             Inches(0.55), Inches(0.35), Inches(11), Inches(0.85),
             size=Pt(34), bold=True, color=WHITE)

    steps = [
        ("01", "Analyse fonctionnelle",
         "Revue des cas d'usage prioritaires avec les équipes métier concernées."),
        ("02", "Évaluation sécurité (PISO)",
         "Audit des dépendances, revue du code source, validation de la politique d'installation HKCU."),
        ("03", "Pilote interne",
         "Déploiement sur poste(s) cible(s), retours utilisateurs, ajustements fonctionnels."),
        ("04", "Mise en production",
         "Packaging NSIS, distribution via canal interne, CI/CD sur tag Git pour les mises à jour."),
    ]

    bw = Inches(2.95)
    bh = Inches(4.4)
    by = Inches(1.8)

    for i, (num, title, text) in enumerate(steps):
        bx = Inches(0.55) + i * (bw + Inches(0.18))
        add_rect(s, bx, by, bw, bh, fill=DARK2)
        add_rect(s, bx, by, bw, Inches(0.05), fill=GREEN)

        add_text(s, num,
                 bx + Inches(0.15), by + Inches(0.12),
                 Inches(0.8), Inches(0.55),
                 size=Pt(28), bold=True, color=GREEN)
        add_text(s, title,
                 bx + Inches(0.15), by + Inches(0.7),
                 bw - Inches(0.25), Inches(0.55),
                 size=Pt(15), bold=True, color=WHITE)
        add_text(s, text,
                 bx + Inches(0.15), by + Inches(1.35),
                 bw - Inches(0.25), bh - Inches(1.45),
                 size=Pt(13), color=MIDGRAY, wrap=True)


# ── Main ──────────────────────────────────────────────────────────────────────

def main():
    prs = new_prs()

    screenshot = r'D:\Developpement\PDF-Equilibrist\docs\présentation\screenshot_v1.png'

    slide_titre(prs)
    slide_apercu(prs, screenshot)
    slide_besoin(prs)
    slide_fonctionnalites(prs)
    slide_utilisateurs(prs)
    slide_technique(prs)
    slide_valeur(prs)
    slide_etapes(prs)

    out = r'D:\Developpement\PDF-Equilibrist\docs\présentation\2026-07-16 - Présentation PDF Equilibrist v2.pptx'
    prs.save(out)
    print(f'Saved: {out}')


if __name__ == '__main__':
    main()
