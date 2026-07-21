from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


BASE_TEMPLATE = Path(r"D:\Developpement\PDF-Equilibrist\dist\PDF-Equilibrist\_internal\pptx\templates\default.pptx")
OUTPUT = Path(r"D:\Developpement\PDF-Equilibrist\release\PDF-Equilibrist_Besoin_Metier.pptx")
ACCENT = RGBColor(107, 191, 78)
DARK = RGBColor(30, 30, 30)
TEXT = RGBColor(240, 240, 240)


def set_text_style(shape, text, size=20, color=TEXT, bold=False, align=PP_ALIGN.LEFT, level=0):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = "Segoe UI"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return run


def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_box = slide.shapes.title
    title_box.text = title
    title_box.text_frame.paragraphs[0].font.size = Pt(26)
    title_box.text_frame.paragraphs[0].font.bold = True
    title_box.text_frame.paragraphs[0].font.color.rgb = ACCENT

    subtitle_box = slide.placeholders[1]
    subtitle_box.text = subtitle
    subtitle_box.text_frame.paragraphs[0].font.size = Pt(14)
    subtitle_box.text_frame.paragraphs[0].font.color.rgb = TEXT


def add_section_slide(prs, title, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    slide.shapes.title.text = title
    if subtitle:
        box = slide.shapes.add_textbox(Inches(0.7), Inches(1.4), Inches(11.5), Inches(0.4))
        set_text_style(box, subtitle, size=14, color=TEXT)


def add_bullets_slide(prs, title, bullets, subtext=None):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    tf.word_wrap = True
    for idx, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(18)
        p.font.name = "Segoe UI"
        p.font.color.rgb = TEXT
        p.bullet = True
    if subtext:
        p = tf.add_paragraph()
        p.text = subtext
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(200, 200, 200)


def add_two_column_slide(prs, title, left_items, right_items):
    slide = prs.slides.add_slide(prs.slide_layouts[3])
    slide.shapes.title.text = title

    left_box = slide.shapes[1]
    right_box = slide.shapes[2]
    left_tf = left_box.text_frame
    right_tf = right_box.text_frame
    left_tf.clear(); right_tf.clear()

    for idx, item in enumerate(left_items):
        p = left_tf.paragraphs[0] if idx == 0 else left_tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(16)
        p.font.name = "Segoe UI"
        p.font.color.rgb = TEXT
        p.bullet = True

    for idx, item in enumerate(right_items):
        p = right_tf.paragraphs[0] if idx == 0 else right_tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(16)
        p.font.name = "Segoe UI"
        p.font.color.rgb = TEXT
        p.bullet = True


def build_deck():
    prs = Presentation(str(BASE_TEMPLATE))

    add_title_slide(
        prs,
        "PDF Equilibrist",
        "Besoin métier et fonctionnalités clés\nPrésentation orientée analyse métier / cadrage fonctionnel",
    )

    add_section_slide(
        prs,
        "1. Le besoin métier couvert",
        "Une solution d’édition, conversion et sécurisation de documents PDF, pensée pour un usage professionnel Windows, simple et robuste.",
    )

    add_bullets_slide(
        prs,
        "Pourquoi ce besoin existe",
        [
            "Les PDF sont largement utilisés pour délivrer des documents finals, mais restent souvent difficiles à corriger ou à enrichir sans outils lourds ou coûteux.",
            "Les équipes métier doivent pouvoir modifier rapidement du texte, ajouter des annotations, préparer des livrables et sécuriser les documents sans sortir du contexte Windows.",
            "Le besoin cible une expérience de traitement PDF complète : visualisation, édition, conversion, impression, protection et contrôle d’accès.",
        ],
    )

    add_two_column_slide(
        prs,
        "2. Utilisateurs et cas d’usage prioritaires",
        [
            "Équipes administratives : correction de texte, mise à jour de documents de référence.",
            "Services de production / documentation : ajout de filigranes, tampons, signatures, annotations.",
            "Direction / coordination : contrôle de la diffusion et sécurisation des documents sensibles.",
        ],
        [
            "Convertisseurs de contenus Office vers PDF et inversement.",
            "Gestion de dossiers PDF multi-pages : découpage, fusion, rotation, inversion.",
            "Préparation de documents pour impression métier ou diffusion interne / externe.",
        ],
    )

    add_bullets_slide(
        prs,
        "3. Fonctionnalités clés attendues",
        [
            "Affichage et navigation dans les PDF, gestion multi-document, vue d’aperçu des miniatures.",
            "Édition de texte, ajout de texte libre, images, filigranes, tampons et signatures.",
            "Annotations métier : surlignage, barré, souligné, zone de texte libre.",
            "Manipulation des pages : rotation, inversion, insertion, split, merge, redimensionnement et recadrage.",
            "Conversion vers/depuis PDF : Office → PDF, PDF → images, Word, Excel, PowerPoint, image → PDF.",
            "Protection : chiffrement AES-256, déchiffrement, gestion des permissions et impression sécurisée.",
        ],
    )

    add_bullets_slide(
        prs,
        "4. Valeur fonctionnelle attendue",
        [
            "Réduction du recours à plusieurs outils disparates pour traiter un même document PDF.",
            "Amélioration de la productivité des équipes en intégrant un flux de travail unique, local et contrôlable.",
            "Renforcement de la conformité et de la traçabilité via protection des documents et gestion des permissions d’usage.",
            "Adaptation à des documents techniques, graphiques et orientés plan / métier, avec rendu imprimable de qualité.",
        ],
    )

    add_bullets_slide(
        prs,
        "5. Positionnement fonctionnel",
        [
            "L’application couvre un besoin d’éditeur PDF desktop robuste, orienté usages bureautiques et documents métiers sensibles.",
            "Elle combine lecture, édition, annotation, conversion et protection dans une seule interface cohérente.",
            "Le périmètre apporte une vision claire du besoin métier et permet de lancer l’analyse fonctionnelle avec une base solide de cas d’usage.",
        ],
    )

    prs.save(str(OUTPUT))
    print(f"Presentation generated: {OUTPUT}")


if __name__ == "__main__":
    build_deck()
